# --- START OF FILE data_ingestion.py ---

from __future__ import annotations
import os
import sys
import json
import uuid
import hashlib
import shutil
from pathlib import Path
from datetime import datetime, timezone
from typing import Iterable, List, Optional, Dict, Any

import fitz  # PyMuPDF for PDF
import docx # python-docx for .docx
import pandas as pd # pandas for .xlsx, .csv
from PIL import Image # Pillow for images
import pytesseract # For OCR on images
from pptx import Presentation # python-pptx for .pptx

from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader # Keep for consistency, but we'll manually handle
from langchain_community.vectorstores import FAISS

from utils.model_loader import ModelLoader
from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException

# We will directly use FileProcessor for saving, so `save_uploaded_files` might not be directly used here anymore
# from utils.file_io import generate_session_id, save_uploaded_files 
from utils.file_io import generate_session_id # Keep this for session ID generation
from utils.document_ops import load_documents, concat_for_analysis, concat_for_comparison # These might need review depending on how they handle new types

# Update SUPPORTED_EXTENSIONS
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".xlsx", ".csv", ".pptx", ".png", ".jpg", ".jpeg"}

# Ensure Tesseract path is set if not in PATH (Windows example)
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


# FAISS Manager (load-or-create)
class FaissManager:
    def __init__(self, index_dir: Path, model_loader: Optional[ModelLoader] = None):
        self.index_dir = Path(index_dir)
        self.index_dir.mkdir(parents=True, exist_ok=True)

        self.meta_path = self.index_dir / "ingested_meta.json"
        self._meta: Dict[str, Any] = {"rows": {}}

        if self.meta_path.exists():
            try:
                self._meta = json.loads(self.meta_path.read_text(encoding="utf-8")) or {"rows": {}}
            except Exception:
                self._meta = {"rows": {}}


        self.model_loader = model_loader or ModelLoader()
        self.emb = self.model_loader.load_embeddings()
        self.vs: Optional[FAISS] = None

    def _exists(self)-> bool:
        return (self.index_dir / "index.faiss").exists() and (self.index_dir / "index.pkl").exists()

    @staticmethod
    def _fingerprint(text: str, md: Dict[str, Any]) -> str:
        src = md.get("source") or md.get("file_path")
        rid = md.get("row_id")
        if src is not None:
            return f"{src}::{'' if rid is None else rid}"
        return hashlib.sha256(text.encode("utf-8")).hexdigest()

    def _save_meta(self):
        self.meta_path.write_text(json.dumps(self._meta, ensure_ascii=False, indent=2), encoding="utf-8")


    def add_documents(self,docs: List[Document]):
        if self.vs is None:
            raise RuntimeError("Call load_or_create() before add_documents_idempotent().")

        new_docs: List[Document] = []

        for d in docs:
            key = self._fingerprint(d.page_content, d.metadata or {})
            if key in self._meta["rows"]:
                continue
            self._meta["rows"][key] = True
            new_docs.append(d)

        if new_docs:
            self.vs.add_documents(new_docs)
            self.vs.save_local(str(self.index_dir))
            self._save_meta()
        return len(new_docs)

    def load_or_create(self,texts:Optional[List[str]]=None, metadatas: Optional[List[dict]] = None):
        if self._exists():
            self.vs = FAISS.load_local(
                str(self.index_dir),
                embeddings=self.emb,
                allow_dangerous_deserialization=True,
            )
            return self.vs
        if not texts:
            raise DocumentPortalException("No existing FAISS index and no data to create one", sys)

        self.vs = FAISS.from_texts(texts=texts, embedding=self.emb, metadatas=metadatas or [])
        self.vs.save_local(str(self.index_dir))
        return self.vs


class ChatIngestor:
    def __init__( self,
        temp_base: str = "data",
        faiss_base: str = "faiss_index",
        use_session_dirs: bool = True,
        session_id: Optional[str] = None,
    ):
        try:
            self.log = CustomLogger().get_logger(__name__)
            self.model_loader = ModelLoader()

            self.use_session = use_session_dirs
            self.session_id = session_id or generate_session_id()

            self.temp_base = Path(temp_base); self.temp_base.mkdir(parents=True, exist_ok=True)
            self.faiss_base = Path(faiss_base); self.faiss_base.mkdir(parents=True, exist_ok=True)

            self.temp_dir = self._resolve_dir(self.temp_base)
            self.faiss_dir = self._resolve_dir(self.faiss_base)

            self.log.info("ChatIngestor initialized",
                          session_id=self.session_id,
                          temp_dir=str(self.temp_dir),
                          faiss_dir=str(self.faiss_dir),
                          sessionized=self.use_session)
        except Exception as e:
            self.log.error("Failed to initialize ChatIngestor", error=str(e))
            raise DocumentPortalException("Initialization error in ChatIngestor", e) from e


    def _resolve_dir(self, base: Path):
        if self.use_session:
            d = base / self.session_id # e.g. "faiss_index/abc123"
            d.mkdir(parents=True, exist_ok=True) # creates dir if not exists
            return d
        return base # fallback: "faiss_index/"

    def _split(self, docs: List[Document], chunk_size=1000, chunk_overlap=200) -> List[Document]:
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        chunks = splitter.split_documents(docs)
        self.log.info("Documents split", chunks=len(chunks), chunk_size=chunk_size, overlap=chunk_overlap)
        return chunks

    def built_retriver( self,
        uploaded_files: Iterable, # This will be an Iterable of FastAPIFileAdapter instances
        *,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        k: int = 5,):
        try:
            # Instantiate FileProcessor to handle saving and reading
            fp = FileProcessor(self.temp_dir, self.session_id)
            
            # Use FileProcessor's save_multiple_files to handle all uploaded files
            # This ensures consistency and leverages FileProcessor's checks.
            paths = fp.save_multiple_files(uploaded_files)
            
            all_documents = []
            for p in paths:
                content = fp.read_file(p)
                if content:
                    all_documents.append(Document(page_content=content, metadata={"source": str(p)}))

            if not all_documents:
                raise ValueError("No valid documents loaded from uploaded files.")

            chunks = self._split(all_documents, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
            fm = FaissManager(self.faiss_dir, self.model_loader)

            texts = [c.page_content for c in chunks]
            metas = [c.metadata for c in chunks]

            try:
                vs = fm.load_or_create(texts=texts, metadatas=metas)
            except Exception: # This bare except might hide issues, consider specific exceptions
                # If load_or_create fails the first time (e.g., index corrupted or not found after checking exists())
                # it's unlikely to succeed on a retry with the same parameters.
                # This 'except' block might be redundant or indicative of a deeper issue if hit.
                # Keeping it for now as it was in your original code, but note this behavior.
                vs = fm.load_or_create(texts=texts, metadatas=metas)

            added = fm.add_documents(chunks)
            self.log.info("FAISS index updated", added=added, index=str(self.faiss_dir))

            return vs.as_retriever(search_type="similarity", search_kwargs={"k": k})

        except Exception as e:
            self.log.error("Failed to build retriever", error=str(e))
            raise DocumentPortalException("Failed to build retriever", e) from e


class FileProcessor:
    """
    Handles saving and reading various document types (PDF, DOCX, TXT, XLSX, CSV, PPTX, Images)
    for analysis and comparison.
    """
    def __init__(self, base_dir: Optional[str] = None, session_id: Optional[str] = None):
        self.log = CustomLogger().get_logger(__name__)
        self.base_dir = Path(base_dir or os.getenv("DATA_STORAGE_PATH", os.path.join(os.getcwd(), "data", "processed_docs")))
        self.session_id = session_id or generate_session_id("session")
        self.session_path = self.base_dir / self.session_id
        self.session_path.mkdir(parents=True, exist_ok=True)
        self.log.info("FileProcessor initialized", session_id=self.session_id, session_path=str(self.session_path))

    def _save_file_content(self, uploaded_file, save_path: Path):
        """Helper to save file content from an uploaded file object."""
        with open(save_path, "wb") as f:
            if hasattr(uploaded_file, "read"):
                f.write(uploaded_file.read())
            else:
                f.write(uploaded_file.getbuffer())

    def save_file(self, uploaded_file) -> Path:
        """Saves a single uploaded file to the session directory."""
        try:
            filename = uploaded_file.name
            file_extension = Path(filename).suffix.lower()

            if file_extension not in SUPPORTED_EXTENSIONS:
                raise ValueError(f"Unsupported file type: {file_extension}. Supported types are: {', '.join(SUPPORTED_EXTENSIONS)}")

            save_path = self.session_path / filename
            self._save_file_content(uploaded_file, save_path)
            self.log.info("File saved successfully", file=filename, save_path=str(save_path), session_id=self.session_id)
            return save_path
        except Exception as e:
            self.log.error("Failed to save file", error=str(e), file_name=uploaded_file.name, session_id=self.session_id)
            raise DocumentPortalException(f"Failed to save file {uploaded_file.name}: {str(e)}", e) from e

    def save_multiple_files(self, uploaded_files: Iterable) -> List[Path]:
        """Saves multiple uploaded files."""
        saved_paths = []
        for uploaded_file in uploaded_files:
            # uploaded_file here is expected to be an object with .name and .read()/.getbuffer()
            # If coming directly from FastAPI, it will be an UploadFile, which FastAPIFileAdapter handles.
            saved_paths.append(self.save_file(uploaded_file))
        return saved_paths

    def read_file(self, file_path: Path) -> Optional[str]:
        """Reads content from various file types."""
        file_extension = file_path.suffix.lower()
        content = None
        try:
            if file_extension == ".pdf":
                text_chunks = []
                with fitz.open(file_path) as doc:
                    if doc.is_encrypted:
                        self.log.warning("Encrypted PDF skipped", file=str(file_path))
                        return None
                    for page_num in range(doc.page_count):
                        page = doc.load_page(page_num)
                        text_chunks.append(f"\n--- Page {page_num + 1} ---\n{page.get_text()}")  # type: ignore
                content = "\n".join(text_chunks)
                self.log.info("PDF read successfully", file=str(file_path), pages=len(text_chunks))

            elif file_extension == ".docx":
                doc = docx.Document(file_path)
                content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
                self.log.info("DOCX read successfully", file=str(file_path))

            elif file_extension == ".txt":
                content = file_path.read_text(encoding="utf-8")
                self.log.info("TXT read successfully", file=str(file_path))

            elif file_extension == ".xlsx":
                # Read all sheets into a dictionary of DataFrames
                dfs = pd.read_excel(file_path, sheet_name=None)
                sheet_contents = []
                for sheet_name, df in dfs.items():
                    # Convert each sheet's DataFrame to a string representation
                    # You might want to customize this for large sheets (e.g., sample rows)
                    sheet_contents.append(f"\n--- Sheet: {sheet_name} ---\n{df.to_csv(index=False)}")
                content = "\n".join(sheet_contents)
                self.log.info("XLSX read successfully", file=str(file_path), sheets=len(dfs))

            elif file_extension == ".csv":
                df = pd.read_csv(file_path)
                content = df.to_csv(index=False) # Convert DataFrame to CSV string
                self.log.info("CSV read successfully", file=str(file_path))

            elif file_extension == ".pptx":
                prs = Presentation(file_path)
                slide_contents = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    if slide_text:
                        slide_contents.append(f"\n--- Slide {i + 1} ---\n{' '.join(slide_text)}")
                content = "\n".join(slide_contents)
                self.log.info("PPTX read successfully", file=str(file_path), slides=len(prs.slides))

            elif file_extension in [".png", ".jpg", ".jpeg"]:
                try:
                    img = Image.open(file_path)
                    content = pytesseract.image_to_string(img)
                    if not content.strip():
                        self.log.warning("No text found in image via OCR", file=str(file_path))
                        return None
                    self.log.info("Image OCR processed successfully", file=str(file_path))
                except pytesseract.TesseractNotFoundError:
                    self.log.error("Tesseract is not installed or not in PATH. Cannot process image files.", file=str(file_path))
                    return None
                except Exception as img_e:
                    self.log.error(f"Error processing image {file_path.name}: {img_e}", file=str(file_path))
                    return None
            else:
                self.log.warning(f"Unsupported file type for reading: {file_extension}", file=str(file_path))
                return None # Explicitly return None for unsupported types

            return content.strip() if content else None

        except Exception as e:
            self.log.error("Failed to read file content", error=str(e), file_path=str(file_path), session_id=self.session_id)
            raise DocumentPortalException(f"Could not process file {file_path.name}: {str(e)}", e) from e

    def combine_documents_in_session(self) -> str:
        """Combines text content of all supported documents in the current session."""
        combined_text_parts = []
        for file_path in sorted(self.session_path.iterdir()):
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                content = self.read_file(file_path)
                if content:
                    combined_text_parts.append(f"Document: {file_path.name}\n{content}")
            else:
                self.log.debug(f"Skipping unsupported file or directory in session: {file_path.name}")
        combined_text = "\n\n".join(combined_text_parts)
        self.log.info("Documents combined from session", count=len(combined_text_parts), session=self.session_id)
        return combined_text

    def clean_old_sessions(self, base_dir: Optional[Path] = None, keep_latest: int = 3):
        """Cleans up old session directories."""
        target_dir = base_dir or self.base_dir
        try:
            sessions = sorted([f for f in target_dir.iterdir() if f.is_dir()], reverse=True)
            for folder in sessions[keep_latest:]:
                shutil.rmtree(folder, ignore_errors=True)
                self.log.info("Old session folder deleted", path=str(folder))
        except Exception as e:
            self.log.error("Error cleaning old sessions", error=str(e))
            raise DocumentPortalException("Error cleaning old sessions", e) from e


# --- END OF FILE data_ingestion.py ---